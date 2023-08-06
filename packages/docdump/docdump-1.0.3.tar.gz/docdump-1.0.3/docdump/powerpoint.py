from pptx import Presentation

class powerpoint_reader:
	def __init__(self):
		pass

	def read(self, path):
		'''
		Opens a powerpoint file, extracts text data

		inputs: str, a path to a PDF file
		outputs: str, a string of all the text in the pdf that could be extracted.
		'''
		prs = Presentation(path)
		text_dump = []

		for slide in prs.slides:
			for shape in slide.shapes:
				
				if hasattr(shape, "text"):
					text_dump.append(shape.text)

		return " ".join(text_dump)
